# VERSION 11.2: Stability & Performance — page_config order, tab-2 columns, safer Anthropic system blocks, throttled streaming, forms, cached client

import os
import io
import json
import time
from datetime import datetime, timedelta
from collections import defaultdict

import streamlit as st
import anthropic
import docx
from pptx import Presentation
from PyPDF2 import PdfReader

# ---- Page config must be FIRST streamlit call ----
st.set_page_config(page_title="SEL Integration Agent", page_icon="🧠", layout="wide")

# -------------------- CONSTANTS --------------------
GRADE_LEVELS = [
    "Kindergarten", "1st Grade", "2nd Grade", "3rd Grade", "4th Grade", "5th Grade",
    "6th Grade", "7th Grade", "8th Grade", "9th Grade", "10th Grade", "11th Grade", "12th Grade"
]
SUBJECTS = ["Science", "History", "English Language Arts", "Mathematics", "Art", "Music"]
COMPETENCIES = {
    "Self-Awareness": ["Identifying Emotions", "Self-Perception", "Recognizing Strengths", "Self-Confidence", "Self-Efficacy"],
    "Self-Management": ["Impulse Control", "Stress Management", "Self-Discipline", "Self-Motivation", "Goal-Setting", "Organizational Skills"],
    "Social Awareness": ["Perspective-Taking", "Empathy", "Appreciating Diversity", "Respect for Others"],
    "Relationship Skills": ["Communication", "Social Engagement", "Building Relationships", "Teamwork", "Conflict Resolution"],
    "Responsible Decision-Making": ["Identifying Problems", "Analyzing Situations", "Solving Problems", "Evaluating", "Reflecting", "Ethical Responsibility"]
}
CASEL_COMPETENCIES = list(COMPETENCIES.keys())

INPUT_COST_PER_MTK = 3.00
OUTPUT_COST_PER_MTK = 15.00
CACHE_WRITE_COST_PER_MTK = 3.75
CACHE_READ_COST_PER_MTK = 0.30

MAX_CALLS_PER_MINUTE = 50
MAX_CALLS_PER_HOUR = 1000

MODEL_NAME = "claude-sonnet-4-5-20250929"


# -------------------- SESSION DEFAULTS --------------------
SESSION_STATE_DEFAULTS = {
    "ai_response": "", "response_title": "", "student_materials": "",
    "differentiation_response": "", "parent_email": "", "scenario": "",
    "conversation_history": [], "training_module": "", "training_scenario": "",
    "training_feedback": "", "check_in_questions": "", "strategy_response": "",
    "total_tokens_used": 0, "total_api_calls": 0,
    "session_start_time": datetime.now(),
    "api_call_times": [],
    "conversation_memory": [],
    "use_streaming": True,
    "estimated_cost": 0.0,
    "screening_data": {},
    "screening_grade": "3rd Grade",
    "screening_num_students": 20,
    "current_student_index": 0,
    "screening_complete": False,
    "screening_interventions": {}
}
for key, default_value in SESSION_STATE_DEFAULTS.items():
    if key not in st.session_state:
        st.session_state[key] = default_value


# -------------------- API CONFIGURATION --------------------
@st.cache_resource
def get_anthropic_client():
    api_key = st.secrets.get("ANTHROPIC_API_KEY") or os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        st.error("🔴 ANTHROPIC_API_KEY not found. Add it to Streamlit Secrets or set env var.")
        st.stop()
    try:
        return anthropic.Anthropic(api_key=api_key)
    except Exception as e:
        st.error(f"🔴 Error initializing Anthropic client: {e}")
        st.stop()

client = get_anthropic_client()


# -------------------- RATE LIMITING --------------------
class RateLimiter:
    @staticmethod
    def check_rate_limit():
        current_time = datetime.now()
        st.session_state.api_call_times = [
            t for t in st.session_state.api_call_times
            if current_time - t < timedelta(hours=1)
        ]
        recent_calls = [
            t for t in st.session_state.api_call_times
            if current_time - t < timedelta(minutes=1)
        ]
        if len(recent_calls) >= MAX_CALLS_PER_MINUTE:
            return False, f"Rate limit exceeded: Maximum {MAX_CALLS_PER_MINUTE} calls per minute"
        if len(st.session_state.api_call_times) >= MAX_CALLS_PER_HOUR:
            return False, f"Rate limit exceeded: Maximum {MAX_CALLS_PER_HOUR} calls per hour"
        return True, "OK"

    @staticmethod
    def record_api_call():
        st.session_state.api_call_times.append(datetime.now())
        st.session_state.total_api_calls += 1


# -------------------- USAGE TRACKING --------------------
class UsageTracker:
    @staticmethod
    def update_usage(input_tokens, output_tokens, cache_creation_tokens=0, cache_read_tokens=0):
        total_tokens = input_tokens + output_tokens
        st.session_state.total_tokens_used += total_tokens
        input_cost = (input_tokens / 1_000_000) * INPUT_COST_PER_MTK
        output_cost = (output_tokens / 1_000_000) * OUTPUT_COST_PER_MTK
        cache_write_cost = (cache_creation_tokens / 1_000_000) * CACHE_WRITE_COST_PER_MTK
        cache_read_cost = (cache_read_tokens / 1_000_000) * CACHE_READ_COST_PER_MTK
        st.session_state.estimated_cost += (input_cost + output_cost + cache_write_cost + cache_read_cost)

    @staticmethod
    def get_usage_summary():
        session_duration = datetime.now() - st.session_state.session_start_time
        hours = session_duration.total_seconds() / 3600
        return {
            "total_calls": st.session_state.total_api_calls,
            "total_tokens": st.session_state.total_tokens_used,
            "estimated_cost": st.session_state.estimated_cost,
            "session_duration": session_duration,
            "calls_per_hour": st.session_state.total_api_calls / hours if hours > 0 else 0
        }


# -------------------- CONVERSATION MEMORY --------------------
class ConversationMemory:
    @staticmethod
    def add_to_memory(role, content, metadata=None):
        memory_entry = {
            "role": role,
            "content": content,
            "timestamp": datetime.now().isoformat(),
            "metadata": metadata or {}
        }
        st.session_state.conversation_memory.append(memory_entry)
        if len(st.session_state.conversation_memory) > 40:
            st.session_state.conversation_memory = st.session_state.conversation_memory[-40:]

    @staticmethod
    def get_relevant_context(current_topic=None, max_messages=10):
        return st.session_state.conversation_memory[-max_messages:]

    @staticmethod
    def format_context_for_prompt():
        if not st.session_state.conversation_memory:
            return ""
        context_parts = ["Previous conversation context:"]
        for entry in st.session_state.conversation_memory[-10:]:
            role = entry['role']
            content = entry['content'][:200]
            context_parts.append(f"{role}: {content}...")
        return "\n".join(context_parts)


# -------------------- HELPERS --------------------
def read_document(uploaded_file):
    if not uploaded_file:
        return ""
    file_extension = os.path.splitext(uploaded_file.name)[1].lower()
    file_bytes = io.BytesIO(uploaded_file.read())
    text_content = ""
    try:
        if file_extension == ".docx":
            doc = docx.Document(file_bytes)
            text_content = "\n".join([para.text for para in doc.paragraphs])
        elif file_extension == ".pptx":
            prs = Presentation(file_bytes)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"
        elif file_extension == ".pdf":
            reader = PdfReader(file_bytes)
            for page in reader.pages:
                t = page.extract_text() or ""
                text_content += t + "\n"
        elif file_extension == ".txt":
            text_content = file_bytes.read().decode("utf-8")
    except Exception as e:
        st.error(f"Error reading file: {e}")
        return ""
    return text_content


def create_docx(text):
    doc = docx.Document()
    doc.add_heading('SEL Integration Plan', 0)
    for line in text.split('\n'):
        if line.startswith('### '):
            doc.add_heading(line.lstrip('### '), level=3)
        elif line.startswith('## '):
            doc.add_heading(line.lstrip('## '), level=2)
        elif line.startswith('# '):
            doc.add_heading(line.lstrip('# '), level=1)
        else:
            doc.add_paragraph(line)
    docx_file = io.BytesIO()
    doc.save(docx_file)
    docx_file.seek(0)
    return docx_file


# -------------------- LLM CALLS --------------------
SYSTEM_PROMPT = """
You are an expert SEL (Social-Emotional Learning) consultant supporting K–12 educators. Your guidance is practical, evidence-based, and grounded in research from CASEL, ASCA, and peer-reviewed educational psychology journals.

Core Directives:
- All primary recommendations MUST follow this exact four-part Markdown format:
  **Overview:** Brief definition or contextual framing (1-2 sentences).
  **Evidence Summary:** What research demonstrates, including study types and key findings (2-4 sentences). Explicitly reference evidence types (e.g., "A meta-analysis of over 200 school-based programs...", "Consistent with developmental research...", "Validated through RCTs...").
  **Implementation Example:** A concrete classroom application with 3-5 actionable steps.
  **Measurement/Outcome:** Observable indicators of success and progress tracking methods (2-3 measurable criteria).
- Link strategies to specific CASEL competencies.
- Maintain a professional, compassionate, data-driven tone.
- Prioritize meta-analyses and systematic reviews over single studies.
"""

def _system_blocks(use_cache: bool):
    blk = {"type": "text", "text": SYSTEM_PROMPT}
    if use_cache:
        blk["cache_control"] = {"type": "ephemeral"}
    return [blk]


def call_claude_streaming(prompt, max_tokens=4096, temperature=1.0, use_cache=True):
    ok, msg = RateLimiter.check_rate_limit()
    if not ok:
        st.error(f"⚠️ {msg}. Please wait a moment.")
        return None
    try:
        response_placeholder = st.empty()
        full_response, buf = "", []
        last = time.time()
        RateLimiter.record_api_call()
        with client.messages.stream(
            model=MODEL_NAME,
            max_tokens=max_tokens,
            temperature=temperature,
            system=_system_blocks(use_cache),
            messages=[{"role": "user", "content": prompt}]
        ) as stream:
            for tok in stream.text_stream:
                buf.append(tok)
                if len(buf) >= 40 or time.time() - last > 0.04:
                    full_response += "".join(buf)
                    buf, last = [], time.time()
                    response_placeholder.markdown(full_response + "▌")
        if buf:
            full_response += "".join(buf)
        response_placeholder.markdown(full_response)
        usage = stream.get_final_message().usage
        UsageTracker.update_usage(
            input_tokens=usage.input_tokens,
            output_tokens=usage.output_tokens,
            cache_creation_tokens=getattr(usage, 'cache_creation_input_tokens', 0),
            cache_read_tokens=getattr(usage, 'cache_read_input_tokens', 0)
        )
        ConversationMemory.add_to_memory("assistant", full_response)
        return full_response
    except anthropic.APIError as e:
        st.error(f"API Error: {e}")
        return None
    except Exception as e:
        st.error(f"Unexpected error: {e}")
        return None


def call_claude(prompt, max_tokens=4096, temperature=1.0, use_cache=True, stream=None):
    should_stream = stream if stream is not None else st.session_state.use_streaming
    if should_stream:
        return call_claude_streaming(prompt, max_tokens, temperature, use_cache)

    ok, msg = RateLimiter.check_rate_limit()
    if not ok:
        st.error(f"⚠️ {msg}. Please wait a moment.")
        return None
    try:
        RateLimiter.record_api_call()
        message = client.messages.create(
            model=MODEL_NAME,
            max_tokens=max_tokens,
            temperature=temperature,
            system=_system_blocks(use_cache),
            messages=[{"role": "user", "content": prompt}]
        )
        response_text = message.content[0].text
        UsageTracker.update_usage(
            input_tokens=message.usage.input_tokens,
            output_tokens=message.usage.output_tokens,
            cache_creation_tokens=getattr(message.usage, 'cache_creation_input_tokens', 0),
            cache_read_tokens=getattr(message.usage, 'cache_read_input_tokens', 0)
        )
        ConversationMemory.add_to_memory("assistant", response_text)
        return response_text
    except anthropic.APIError as e:
        st.error(f"API Error: {e}")
        return None
    except Exception as e:
        st.error(f"Unexpected error: {e}")
        return None


# -------------------- PROMPTS --------------------
def get_analysis_prompt(lesson_plan_text, standard="", competency="", skill=""):
    focus_instruction = ""
    if competency and skill:
        focus_instruction = f"The user has requested specific focus on the CASEL competency of **{competency}**, emphasizing the skill of **{skill}**. Prioritize this focus in your analysis."
    standard_instruction = ""
    if standard and standard.strip():
        standard_instruction = f"All suggestions must align with this educational standard: '{standard.strip()}'."
    context = ConversationMemory.format_context_for_prompt()
    context_section = f"\n\n{context}\n" if context else ""
    return f"""{context_section}

An educator has submitted this lesson plan for SEL integration analysis:

**Lesson Plan:**
---
{lesson_plan_text}
---

**Task:**
1. Analyze the lesson to identify the strongest opportunity for SEL integration.
2. Provide ONE comprehensive, high-impact SEL strategy recommendation.
3. Follow the mandatory four-part format from the system prompt.

{focus_instruction}
{standard_instruction}
"""


def get_creation_prompt(grade_level, subject, topic, competency="", skill=""):
    focus_instruction = ""
    if competency and skill:
        focus_instruction = f"The lesson's primary SEL focus must be **{competency}**, specifically developing **{skill}**."
    context = ConversationMemory.format_context_for_prompt()
    context_section = f"\n\n{context}\n" if context else ""
    return f"""{context_section}

Create a complete, SEL-integrated lesson plan with these specifications:
- **Grade Level:** {grade_level}
- **Subject:** {subject}
- **Topic:** {topic}
- **SEL Focus:** {focus_instruction if focus_instruction else "Balanced approach across CASEL competencies"}

**Requirements:**
1. Start with a "Pedagogical Rationale" (2-3 sentences) explaining the evidence behind your primary SEL activity.
2. Generate a complete lesson plan in Markdown with:
   - Learning Objectives (Content + SEL, observable behaviors)
   - Key Vocabulary (Content + SEL terms)
   - Materials list
   - Lesson Sequence (Hook → I Do → We Do → You Do → Assessment → Closing)
   - Detailed SEL Alignment section

Follow an "I Do, We Do, You Do" instructional model.
"""


def get_strategy_prompt(situation):
    context = ConversationMemory.format_context_for_prompt()
    context_section = f"\n\n{context}\n" if context else ""
    return f"""{context_section}

A teacher needs an immediate, evidence-based strategy for this situation:

**Situation:** "{situation}"

**Task:**
Provide ONE quick, actionable strategy using this format:
- **Strategy Name:** (e.g., "Mindful Minute")
- **Evidence Rationale:** (1-2 sentences on research basis)
- **Actionable Steps:** (2-3 immediate steps)
- **Expected Outcome:** (1 sentence on observable results)
"""


def get_student_materials_prompt(lesson_plan_output):
    return f"""You are an instructional designer. Based on this lesson plan, create student-facing materials in Markdown format:

**Lesson Plan:**
---
{lesson_plan_output}
---

**Generate:**
### 🎟️ Exit Ticket
(2-3 reflective questions)

### 🗣️ Think-Pair-Share Prompts
(2-3 discussion questions)

### ✍️ Journal Starters
(2-3 reflective prompts)

### 📄 Practice Worksheet
(Simple printable worksheet/graphic organizer)
"""


def get_differentiation_prompt(lesson_plan_output):
    return f"""You are an expert in instructional differentiation. Based on this lesson, provide evidence-based strategies in Markdown:

**Lesson Plan:**
---
{lesson_plan_output}
---

**Structure:**
### 📉 Scaffold Support (Struggling Learners)
### ⬆️ Extension Activities (Advanced Learners)
### 🌐 ELL Support
"""


def get_scenario_prompt(competency, skill, grade_level):
    return f"""Generate a brief, relatable school scenario for a {grade_level} student requiring use of the SEL competency **{competency}** (skill: **{skill}**).

Present in second person ('You are...'), ending with a question. Keep it to one paragraph.
"""


def get_feedback_prompt(scenario, history):
    formatted_history = "\n".join([f"- {entry['role']}: {entry['content']}" for entry in history])
    return f"""You are a supportive SEL coach using a Socratic approach.

**Scenario:** {scenario}

**Conversation History:**
{formatted_history}

Ask ONE reflective question to deepen the student's thinking. Do not give advice. Keep it brief.
"""


def get_training_prompt(competency):
    return f"""Create a professional development module on **{competency}** grounded in CASEL and evidence-based practices.

**Structure:**
## 🧠 Understanding {competency}
(Definition and importance, citing CASEL)

## 🛠️ Evidence-Based Classroom Strategies
For 2-3 key skills, provide:
### Skill: [Name]
* **The Strategy:** (Practical approach)
* **Evidence Base:** (Research summary)
* **Implementation Example:** (Step-by-step)
"""


def get_training_scenario_prompt(competency, training_module_text):
    return f"""Create a brief, challenging classroom scenario to help a teacher practice **{competency}**.

End with an open-ended question. Generate ONLY the scenario and question.
"""


def get_training_feedback_prompt(competency, scenario, teacher_response):
    return f"""You are a supportive SEL coach. The teacher is practicing **{competency}**.

**Scenario:** {scenario}
**Teacher's Response:** {teacher_response}

Provide constructive feedback: affirm a positive aspect, then ask one reflective question.
"""


def get_check_in_prompt(grade_level, tone):
    return f"""Generate 3-4 creative, age-appropriate morning check-in questions for a **{grade_level}** class with a **{tone}** tone.

Format as a numbered list.
"""


def get_parent_email_prompt(lesson_plan):
    return f"""Draft a professional, strengths-based email to parents based on this lesson plan:

**Lesson Plan:**
---
{lesson_plan}
---

**Structure:**
1. Subject Line (clear, informative)
2. What We're Learning (main SEL skill)
3. How We Practiced (brief activity description)
4. Connection at Home (simple conversation starter)
"""


def clear_generated_content():
    keys_to_clear = [
        "ai_response", "response_title", "student_materials",
        "differentiation_response", "parent_email", "scenario",
        "conversation_history", "training_module", "training_scenario",
        "training_feedback", "check_in_questions", "strategy_response"
    ]
    for key in keys_to_clear:
        if key in st.session_state:
            st.session_state[key] = ""
    if "conversation_history" in st.session_state:
        st.session_state.conversation_history = []


# -------------------- SEL SCREENER --------------------
def get_screener_questions(grade_level):
    grade_lower = grade_level.lower()
    if "kinder" in grade_lower or grade_level == "K":
        grade_num = 0
    elif "1st" in grade_lower:
        grade_num = 1
    elif "2nd" in grade_lower:
        grade_num = 2
    elif "3rd" in grade_lower:
        grade_num = 3
    elif "4th" in grade_lower:
        grade_num = 4
    elif "5th" in grade_lower:
        grade_num = 5
    else:
        grade_num = int(''.join(filter(str.isdigit, grade_level))) if any(c.isdigit() for c in grade_level) else 3

    if grade_num <= 1:
        return [
            {"emoji": "😊", "text": "Names feelings like happy, sad, or mad", "competency": "Self-Awareness"},
            {"emoji": "🎯", "text": "Can calm down with adult help", "competency": "Self-Management"},
            {"emoji": "👥", "text": "Is kind to friends", "competency": "Social Awareness"},
            {"emoji": "🤝", "text": "Takes turns and shares toys", "competency": "Relationship Skills"},
            {"emoji": "💭", "text": "Follows class rules", "competency": "Decision-Making"}
        ]
    elif grade_num == 2:
        return [
            {"emoji": "😊", "text": "Recognizes and talks about their feelings", "competency": "Self-Awareness"},
            {"emoji": "🎯", "text": "Uses calming strategies when upset (like deep breaths)", "competency": "Self-Management"},
            {"emoji": "👥", "text": "Shows care for others' feelings", "competency": "Social Awareness"},
            {"emoji": "🤝", "text": "Works well in small groups", "competency": "Relationship Skills"},
            {"emoji": "💭", "text": "Thinks before acting", "competency": "Decision-Making"}
        ]
    else:
        return [
            {"emoji": "😊", "text": "Identifies own emotions and what causes them", "competency": "Self-Awareness"},
            {"emoji": "🎯", "text": "Manages frustration and stays calm independently", "competency": "Self-Management"},
            {"emoji": "👥", "text": "Shows empathy and respects different perspectives", "competency": "Social Awareness"},
            {"emoji": "🤝", "text": "Communicates needs and resolves conflicts peacefully", "competency": "Relationship Skills"},
            {"emoji": "💭", "text": "Makes responsible, thoughtful decisions", "competency": "Decision-Making"}
        ]


def calculate_screening_results():
    if not st.session_state.screening_data:
        return None
    results = {
        "total_students": len(st.session_state.screening_data),
        "students": {},
        "class_averages": {},
        "risk_levels": {"priority": [], "monitor": [], "on_track": []}
    }
    for student_id, scores in st.session_state.screening_data.items():
        avg_score = sum(scores) / len(scores)
        student_results = {
            "scores": scores,
            "average": avg_score,
            "risk_level": "priority" if avg_score < 2.0 else ("monitor" if avg_score < 2.5 else "on_track")
        }
        results["students"][student_id] = student_results
        results["risk_levels"][student_results["risk_level"]].append(student_id)
    competencies = ["Self-Awareness", "Self-Management", "Social Awareness", "Relationship Skills", "Decision-Making"]
    for i, comp in enumerate(competencies):
        scores_for_comp = [scores[i] for scores in st.session_state.screening_data.values()]
        results["class_averages"][comp] = sum(scores_for_comp) / len(scores_for_comp)
    return results


def get_intervention_prompt(student_id, student_results, grade_level):
    scores = student_results["scores"]
    avg = student_results["average"]
    competencies = ["Self-Awareness", "Self-Management", "Social Awareness", "Relationship Skills", "Decision-Making"]
    concerns = []
    strengths = []
    for i, comp in enumerate(competencies):
        if scores[i] < 2.5:
            concerns.append(f"{comp} (score: {scores[i]}/4)")
        elif scores[i] >= 3.0:
            strengths.append(comp)
    return f"""You are an SEL intervention specialist. A {grade_level} student needs support.

**Assessment Results:**
- Overall Average: {avg:.1f}/4.0
- Risk Level: {student_results["risk_level"].title()}

**Areas of Concern:**
{chr(10).join(f"- {c}" for c in concerns) if concerns else "None - student is on track"}

**Strengths:**
{chr(10).join(f"- {s}" for s in strengths) if strengths else "Developing in all areas"}

**Your Task:**
Provide 3-4 specific, actionable Tier 2 interventions for this student. Format as:

**Primary Focus:** [The most critical area to address]

**Recommended Interventions:**

1. **[Intervention Name]**
   - What: [Brief description]
   - How: [2-3 concrete steps]
   - Timeline: [How often/how long]

2. **[Intervention Name]**
   - What: [Brief description]
   - How: [2-3 concrete steps]
   - Timeline: [How often/how long]

3. **[Intervention Name]**
   - What: [Brief description]
   - How: [2-3 concrete steps]
   - Timeline: [How often/how long]

**Progress Monitoring:**
- Check progress in: [timeframe]
- Look for: [specific behavioral changes]

**Family Communication:**
[2-3 sentence summary to share with parents about what we're working on]

Keep it practical, evidence-based, and feasible for a busy classroom teacher.
"""


def get_class_strategies_prompt(results, grade_level):
    class_avgs = results["class_averages"]
    lowest_comp = min(class_avgs, key=class_avgs.get)
    lowest_score = class_avgs[lowest_comp]
    on_track_pct = (len(results["risk_levels"]["on_track"]) / results["total_students"]) * 100
    return f"""You are an SEL curriculum specialist. Analyze this {grade_level} class screening data:

**Class Overview:**
- Total Students: {results["total_students"]}
- On Track: {len(results["risk_levels"]["on_track"])} students ({on_track_pct:.0f}%)
- Need Monitoring: {len(results["risk_levels"]["monitor"])} students
- Priority Support: {len(results["risk_levels"]["priority"])} students

**Class Competency Averages:**
{chr(10).join(f"- {comp}: {score:.1f}/4.0" for comp, score in class_avgs.items())}

**Lowest Area:** {lowest_comp} ({lowest_score:.1f}/4.0)

**Your Task:**
Provide whole-class strategies to strengthen the lowest area. Format as:

**Class Need:** {lowest_comp}

**Whole-Class Strategies:**

1. **[Strategy Name]**
   - What: [Brief description]
   - When: [How to incorporate into daily schedule]
   - Materials: [What's needed]

2. **[Strategy Name]**
   - What: [Brief description]
   - When: [How to incorporate into daily schedule]
   - Materials: [What's needed]

3. **[Strategy Name]**
   - What: [Brief description]
   - When: [How to incorporate into daily schedule]
   - Materials: [What's needed]

**Quick Wins:** [2-3 simple things teacher can start tomorrow]

**Resources:** [Specific curricula, books, or websites that align with this focus]

Keep strategies evidence-based, practical, and engaging for {grade_level} students.
"""

def save_screening_data():
    if not st.session_state.screening_data:
        return None
    results = calculate_screening_results()
    if not results:
        return None
    data = {
        "date": datetime.now().isoformat(),
        "grade": st.session_state.screening_grade,
        "num_students": st.session_state.screening_num_students,
        "screening_data": st.session_state.screening_data,
        "results": results,
        "interventions": st.session_state.screening_interventions
    }
    return json.dumps(data, indent=2)


def load_screening_data(uploaded_file):
    try:
        data = json.loads(uploaded_file.read().decode('utf-8'))
        st.session_state.screening_grade = data.get("grade", "3rd Grade")
        st.session_state.screening_num_students = data.get("num_students", 20)
        st.session_state.screening_data = data.get("screening_data", {})
        st.session_state.screening_interventions = data.get("interventions", {})
        st.session_state.screening_complete = bool(st.session_state.screening_data)
        st.session_state.current_student_index = len(st.session_state.screening_data)
        return True
    except Exception as e:
        st.error(f"Error loading screening data: {e}")
        return False


def create_comprehensive_report():
    results = calculate_screening_results()
    if not results:
        return None
    report_parts = []
    report_parts.append(f"# SEL SCREENING REPORT")
    report_parts.append(f"**Grade:** {st.session_state.screening_grade}")
    report_parts.append(f"**Date:** {datetime.now().strftime('%B %d, %Y')}")
    report_parts.append(f"**Total Students:** {results['total_students']}")
    report_parts.append("\n---\n")
    report_parts.append("## CLASS OVERVIEW")
    on_track_pct = (len(results['risk_levels']['on_track']) / results['total_students']) * 100
    monitor_pct = (len(results['risk_levels']['monitor']) / results['total_students']) * 100
    priority_pct = (len(results['risk_levels']['priority']) / results['total_students']) * 100
    report_parts.append(f"- **On Track:** {len(results['risk_levels']['on_track'])} students ({on_track_pct:.0f}%)")
    report_parts.append(f"- **Monitor:** {len(results['risk_levels']['monitor'])} students ({monitor_pct:.0f}%)")
    report_parts.append(f"- **Priority Support:** {len(results['risk_levels']['priority'])} students ({priority_pct:.0f}%)")
    report_parts.append("")
    report_parts.append("## COMPETENCY AVERAGES")
    for comp, avg in results['class_averages'].items():
        status = "✓ Strong" if avg >= 3.0 else ("⚠ Developing" if avg >= 2.5 else "⚡ Needs Focus")
        report_parts.append(f"- **{comp}:** {avg:.1f}/4.0 ({status})")
    report_parts.append("\n---\n")
    if "class" in st.session_state.screening_interventions:
        report_parts.append("## WHOLE-CLASS STRATEGIES")
        report_parts.append(st.session_state.screening_interventions["class"])
        report_parts.append("\n---\n")
    if results["risk_levels"]["priority"] or results["risk_levels"]["monitor"]:
        report_parts.append("## INDIVIDUAL STUDENT INTERVENTION PLANS")
        competencies = ["Self-Awareness", "Self-Management", "Social Awareness", "Relationship Skills", "Decision-Making"]
        if results["risk_levels"]["priority"]:
            report_parts.append("\n### Priority Support Students")
            for student_id in results["risk_levels"]["priority"]:
                student_data = results["students"][student_id]
                report_parts.append(f"\n#### {student_id}")
                report_parts.append(f"**Average Score:** {student_data['average']:.1f}/4.0")
                report_parts.append("\n**Individual Scores:**")
                for i, comp in enumerate(competencies):
                    score = student_data["scores"][i]
                    report_parts.append(f"- {comp}: {score}/4")
                if student_id in st.session_state.screening_interventions:
                    report_parts.append("\n**Intervention Plan:**")
                    report_parts.append(st.session_state.screening_interventions[student_id])
                report_parts.append("")
        if results["risk_levels"]["monitor"]:
            report_parts.append("\n### Monitor Students")
            for student_id in results["risk_levels"]["monitor"]:
                student_data = results["students"][student_id]
                report_parts.append(f"\n#### {student_id}")
                report_parts.append(f"**Average Score:** {student_data['average']:.1f}/4.0")
                report_parts.append("\n**Individual Scores:**")
                for i, comp in enumerate(competencies):
                    score = student_data["scores"][i]
                    report_parts.append(f"- {comp}: {score}/4")
                if student_id in st.session_state.screening_interventions:
                    report_parts.append("\n**Intervention Plan:**")
                    report_parts.append(st.session_state.screening_interventions[student_id])
                report_parts.append("")
    if results["risk_levels"]["on_track"]:
        report_parts.append("\n---\n")
        report_parts.append("## STUDENTS ON TRACK")
        for student_id in results["risk_levels"]["on_track"]:
            avg = results["students"][student_id]["average"]
            report_parts.append(f"- {student_id}: {avg:.1f}/4.0")
    return "\n".join(report_parts)


# -------------------- UI: SIDEBAR --------------------
def is_admin():
    return False  # placeholder

with st.sidebar:
    user_is_admin = is_admin()
    st.header("⚙️ Settings" if not user_is_admin else "⚙️ Admin Dashboard")
    st.markdown("---")
    st.session_state.use_streaming = st.checkbox(
        "Enable Streaming Responses",
        value=st.session_state.use_streaming,
        help="Show responses in real-time as they're generated"
    )
    st.markdown("---")
    st.subheader("🧠 Conversation Memory")
    memory_count = len(st.session_state.conversation_memory)
    st.caption(f"Messages stored: {memory_count}")
    if st.button("Clear Memory", help="Start fresh with a new conversation"):
        st.session_state.conversation_memory = []
        st.success("Memory cleared!")
        st.rerun()
    st.markdown("---")
    duration = datetime.now() - st.session_state.session_start_time
    minutes = int(duration.total_seconds() // 60)
    st.caption(f"⏱️ Session: {minutes} minutes")


# -------------------- UI: MAIN --------------------
st.title("🧠 SEL Integration Agent")
st.markdown("*Powered by Claude Sonnet 4.5 - Your AI instructional coach for Social-Emotional Learning*")

tab_list = ["Analyze Existing Lesson", "Create New Lesson", "🧑‍🎓 Student Scenarios", "👩‍🏫 Teacher SEL Training", "☀️ Morning Check-in", "🆘 Strategy Finder", "📊 SEL Screener"]
tab1, tab2, tab3, tab4, tab5, tab6, tab7 = st.tabs(tab_list)

# ---- TAB 1: Analyze Existing Lesson (wrapped in form) ----
with tab1:
    st.header("Analyze an Existing Lesson Plan")
    if st.button("🗑️ Clear This Tab", key="clear_tab1"):
        clear_generated_content()
        st.success("Tab cleared!")
        st.rerun()

    with st.form("analyze_form"):
        st.info("Upload or paste a lesson plan. Get one high-impact, evidence-based SEL integration strategy.")
        st.markdown("**Optional: Add a Specific SEL Focus**")
        col1a, col2a = st.columns(2)
        with col1a:
            analyze_competency = st.selectbox("Select a CASEL Competency", options=CASEL_COMPETENCIES, index=None, placeholder="Choose a competency...", key="analyze_comp")
        with col2a:
            analyze_skill = st.selectbox("Select a Focused Skill",
                                         options=COMPETENCIES[analyze_competency] if analyze_competency else [],
                                         index=None, placeholder="Choose a skill...",
                                         key="analyze_skill", disabled=not bool(analyze_competency))
        st.markdown("---")
        uploaded_file = st.file_uploader("Upload a .txt, .docx, .pptx, or .pdf file", type=["txt", "docx", "pptx", "pdf"])
        lesson_text_paste = st.text_area("Or paste the full text of your lesson plan here.", height=200)
        standard_input = st.text_area("(Optional) Paste educational standard(s) here.", placeholder="e.g., CCSS.ELA-LITERACY.RL.5.2", height=100)
        submitted_analyze = st.form_submit_button("🚀 Generate SEL Suggestions")

    if submitted_analyze:
        lesson_content = read_document(uploaded_file) if uploaded_file else (lesson_text_paste or "")
        if not lesson_content.strip():
            st.warning("Please upload or paste a lesson plan to begin.")
        else:
            with st.spinner("🤖 Analyzing lesson with Claude Sonnet 4.5..."):
                clear_generated_content()
                ConversationMemory.add_to_memory("user", f"Analyze lesson plan (competency: {analyze_competency}, skill: {analyze_skill})", {"type": "lesson_analysis"})
                prompt = get_analysis_prompt(lesson_content, standard_input, analyze_competency, analyze_skill)
                response = call_claude(prompt)
                if response:
                    st.session_state.ai_response = response
                    st.session_state.response_title = "Evidence-Based SEL Recommendation"

# ---- TAB 2: Create New Lesson (fixed column scope) ----
with tab2:
    st.header("Create a New, SEL-Integrated Lesson")
    if st.button("🗑️ Clear This Tab", key="clear_tab2"):
        clear_generated_content()
        st.success("Tab cleared!")
        st.rerun()

    st.info("Fill in the details to generate a new lesson plan from scratch.")
    st.markdown("**Optional: Add a Specific SEL Focus**")
    col1c, col2c = st.columns(2)
    with col1c:
        create_competency = st.selectbox("Select a CASEL Competency", options=CASEL_COMPETENCIES, index=None, placeholder="Choose a competency...", key="create_comp")
    with col2c:
        create_skill = st.selectbox("Select a Focused Skill",
                                    options=COMPETENCIES[create_competency] if create_competency else [],
                                    index=None, placeholder="Choose a skill...", key="create_skill",
                                    disabled=not bool(create_competency))
    st.markdown("---")
    with st.form("create_form"):
        create_grade = st.selectbox("Grade Level", options=GRADE_LEVELS, index=0)
        create_subject = st.selectbox("Subject", options=SUBJECTS, index=0)
        create_topic = st.text_area("Lesson Topic or Objective", "The causes and effects of the American Revolution.")
        submitted = st.form_submit_button("✨ Create SEL Lesson Plan")

    if submitted:
        with st.spinner("🛠️ Building your lesson plan with Claude Sonnet 4.5..."):
            clear_generated_content()
            ConversationMemory.add_to_memory("user", f"Create lesson: {create_topic} ({create_grade}, {create_subject})", {"type": "lesson_creation"})
            prompt = get_creation_prompt(create_grade, create_subject, create_topic, create_competency, create_skill)
            response = call_claude(prompt)
            if response:
                st.session_state.ai_response = response
                st.session_state.response_title = "Your New SEL-Integrated Lesson Plan"

# ---- TAB 3: Student Scenarios ----
with tab3:
    st.header("Interactive SEL Scenarios")
    if st.button("🗑️ Clear This Tab", key="clear_tab3"):
        st.session_state.scenario = ""
        st.session_state.conversation_history = []
        st.success("Scenario cleared!")
        st.rerun()

    st.info("Select a competency and skill to generate a practice scenario.")
    col1b, col2b, col3b = st.columns(3)
    with col1b:
        scenario_competency = st.selectbox("Select a CASEL Competency", options=CASEL_COMPETENCIES, index=3, key="scenario_comp")
    with col2b:
        scenario_skill = st.selectbox("Select a Focused Skill", options=COMPETENCIES[scenario_competency], index=0, key="scenario_skill")
    with col3b:
        scenario_grade = st.selectbox("Select a Grade Level", options=GRADE_LEVELS, key="scenario_grade")

    if st.button("🎬 Generate New Scenario"):
        with st.spinner("Writing a scenario..."):
            prompt = get_scenario_prompt(scenario_competency, scenario_skill, scenario_grade)
            response = call_claude(prompt, max_tokens=1024, stream=False)
            if response:
                st.session_state.scenario = response
                st.session_state.conversation_history = []

    if st.session_state.scenario:
        st.markdown("---")
        st.markdown(f"**Scenario:** {st.session_state.scenario}")
        for entry in st.session_state.conversation_history:
            if entry['role'] == 'Student':
                st.markdown(f"> **You:** {entry['content']}")
            else:
                st.markdown(f"**Coach:** {entry['content']}")
        student_response = st.text_input("What would you do or say?", key="student_response_input")
        if st.button("💬 Submit Response"):
            if student_response:
                st.session_state.conversation_history.append({"role": "Student", "content": student_response})
                with st.spinner("Coach is thinking..."):
                    feedback_prompt = get_feedback_prompt(st.session_state.scenario, st.session_state.conversation_history)
                    response = call_claude(feedback_prompt, max_tokens=1024, stream=False)
                    if response:
                        st.session_state.conversation_history.append({"role": "Coach", "content": response})
                        st.rerun()

# ---- TAB 4: Teacher Training ----
with tab4:
    st.header("👩‍🏫 Teacher SEL Training")
    if st.button("🗑️ Clear This Tab", key="clear_tab4"):
        st.session_state.training_module = ""
        st.session_state.training_scenario = ""
        st.session_state.training_feedback = ""
        st.success("Training cleared!")
        st.rerun()

    st.info("Select a competency to begin an in-depth training module.")
    training_competency = st.selectbox("Select a CASEL Competency to learn about", options=CASEL_COMPETENCIES, index=None, placeholder="Choose a competency...", key="training_comp_select")
    if st.button("🎓 Start Training Module"):
        if training_competency:
            with st.spinner("Preparing your training module..."):
                prompt = get_training_prompt(training_competency)
                response = call_claude(prompt)
                if response:
                    st.session_state.training_module = response
                    st.session_state.training_scenario = ""
                    st.session_state.training_feedback = ""
        else:
            st.warning("Please select a competency to begin.")

    if st.session_state.training_module:
        st.markdown("---")
        st.markdown(st.session_state.training_module)
        st.markdown("---")
        st.subheader("🎬 Let's Try It Out")
        if st.button("Generate a Practice Scenario"):
            with st.spinner("Creating a classroom scenario..."):
                prompt = get_training_scenario_prompt(training_competency, st.session_state.training_module)
                response = call_claude(prompt, max_tokens=1024, stream=False)
                if response:
                    st.session_state.training_scenario = response
                    st.session_state.training_feedback = ""
        if st.session_state.training_scenario:
            st.info(st.session_state.training_scenario)
            teacher_response = st.text_area("How would you respond to this scenario?", key="teacher_response_area")
            if st.button("Get Feedback"):
                if teacher_response:
                    with st.spinner("Your coach is reviewing your response..."):
                        prompt = get_training_feedback_prompt(training_competency, st.session_state.training_scenario, teacher_response)
                        response = call_claude(prompt, max_tokens=1024, stream=False)
                        if response:
                            st.session_state.training_feedback = response
                else:
                    st.warning("Please enter your response above.")
            if st.session_state.training_feedback:
                st.markdown("---")
                st.markdown("#### Coach's Feedback")
                st.success(st.session_state.training_feedback)

# ---- TAB 5: Morning Check-in ----
with tab5:
    st.header("☀️ SEL Morning Check-in")
    if st.button("🗑️ Clear This Tab", key="clear_tab5"):
        st.session_state.check_in_questions = ""
        st.success("Questions cleared!")
        st.rerun()

    st.info("Generate creative questions for your morning meeting or class check-in.")
    col1d, col2d = st.columns(2)
    with col1d:
        check_in_grade = st.selectbox("Select a Grade Level", options=GRADE_LEVELS, key="check_in_grade")
    with col2d:
        check_in_tone = st.selectbox("Select a Tone", options=["Calm", "Energetic", "Reflective", "Fun", "Serious"], key="check_in_tone")
    if st.button("❓ Generate Questions"):
        with st.spinner("Coming up with some good questions..."):
            prompt = get_check_in_prompt(check_in_grade, check_in_tone)
            response = call_claude(prompt, max_tokens=1024, stream=False)
            if response:
                st.session_state.check_in_questions = response
    if st.session_state.check_in_questions:
        st.markdown("---")
        st.markdown(st.session_state.check_in_questions)

# ---- TAB 6: Strategy Finder (wrapped in form) ----
with tab6:
    st.header("🆘 On-Demand Strategy Finder")
    if st.button("🗑️ Clear This Tab", key="clear_tab6"):
        st.session_state.strategy_response = ""
        st.success("Strategy cleared!")
        st.rerun()

    st.info("Describe a classroom situation to get immediate, actionable SEL strategies.")
    with st.form("strategy_form"):
        situation = st.text_area("Describe the situation in your classroom:", placeholder="e.g., 'Two students are arguing over a shared resource' or 'My class is very unfocused after lunch.'", height=150)
        submitted_strategy = st.form_submit_button("💡 Find a Strategy")

    if submitted_strategy:
        if situation and situation.strip():
            with st.spinner("Finding effective strategies..."):
                prompt = get_strategy_prompt(situation)
                response = call_claude(prompt, max_tokens=2048)
                if response:
                    st.session_state.strategy_response = response
        else:
            st.warning("Please describe the situation to get a strategy.")
    if st.session_state.strategy_response:
        st.markdown("---")
        st.markdown(st.session_state.strategy_response)

# ---- TAB 7: SEL Screener ----
with tab7:
    st.header("📊 Quick SEL Screener")
    if st.button("🗑️ Reset Screener", key="clear_tab7"):
        st.session_state.screening_data = {}
        st.session_state.screening_interventions = {}
        st.session_state.current_student_index = 0
        st.session_state.screening_complete = False
        st.success("Screener reset!")
        st.rerun()

    st.info("Screen your class in under an hour. Identify students who need support and get AI-powered intervention plans.")
    col_up, col_down = st.columns(2)
    with col_up:
        uploaded_file = st.file_uploader("📁 Load Previous Screening", type=["json"], key="screener_upload")
        if uploaded_file:
            if load_screening_data(uploaded_file):
                st.success("✓ Screening data loaded successfully!")
    with col_down:
        if st.session_state.screening_complete:
            screening_json = save_screening_data()
            if screening_json:
                st.download_button(label="💾 Save Screening Data", data=screening_json, file_name=f"sel_screening_{datetime.now().strftime('%Y%m%d_%H%M')}.json", mime="application/json")

    st.markdown("---")
    if not st.session_state.screening_complete:
        st.subheader("📝 Set Up Your Screening")
        col1, col2 = st.columns(2)
        with col1:
            st.session_state.screening_grade = st.selectbox("Grade Level", options=GRADE_LEVELS[:6],
                                                            index=GRADE_LEVELS[:6].index(st.session_state.screening_grade) if st.session_state.screening_grade in GRADE_LEVELS[:6] else 3,
                                                            key="screener_grade_select")
        with col2:
            st.session_state.screening_num_students = st.number_input("Number of Students", min_value=1, max_value=35,
                                                                      value=st.session_state.screening_num_students, key="screener_num_students")
        if st.button("🚀 Start Screening", type="primary"):
            st.session_state.screening_data = {}
            st.session_state.current_student_index = 0
            st.session_state.screening_complete = False
            st.rerun()

        if st.session_state.current_student_index < st.session_state.screening_num_students:
            st.markdown("---")
            st.subheader(f"Student {st.session_state.current_student_index + 1} of {st.session_state.screening_num_students}")
            default_student_id = f"Student {st.session_state.current_student_index + 1}"
            existing_student_id = None
            existing_ratings = None
            student_ids = list(st.session_state.screening_data.keys())
            if st.session_state.current_student_index < len(student_ids):
                existing_student_id = student_ids[st.session_state.current_student_index]
                existing_ratings = st.session_state.screening_data[existing_student_id]
            student_id = st.text_input("Student ID or Initials (optional)",
                                       value=existing_student_id if existing_student_id else default_student_id,
                                       key=f"student_id_{st.session_state.current_student_index}")
            st.markdown("**Rate what you observe** (1 = Concern, 2 = Developing, 3 = On Track, 4 = Strong)")
            questions = get_screener_questions(st.session_state.screening_grade)
            st.caption(f"📚 Questions for: {st.session_state.screening_grade}")
            ratings = []
            for i, q in enumerate(questions):
                st.markdown(f"**{q['emoji']} {q['text']}**")
                if existing_ratings and i < len(existing_ratings):
                    default_index = existing_ratings[i] - 1
                else:
                    default_index = 2
                rating = st.radio(
                    f"Rating for question {i+1}",
                    options=[1, 2, 3, 4],
                    format_func=lambda x: ["1 - Concern", "2 - Developing", "3 - On Track", "4 - Strong"][x-1],
                    horizontal=True, index=default_index,
                    key=f"rating_{st.session_state.current_student_index}_{i}_{st.session_state.screening_grade}",
                    label_visibility="collapsed"
                )
                ratings.append(rating)
                st.markdown("")

            col_prev, col_next = st.columns(2)
            with col_prev:
                if st.session_state.current_student_index > 0:
                    if st.button("⬅️ Previous Student"):
                        st.session_state.screening_data[student_id] = ratings
                        st.session_state.current_student_index -= 1
                        st.rerun()
            with col_next:
                button_label = "Next Student ➡️" if st.session_state.current_student_index < st.session_state.screening_num_students - 1 else "✅ Complete Screening"
                if st.button(button_label, type="primary"):
                    st.session_state.screening_data[student_id] = ratings
                    if st.session_state.current_student_index < st.session_state.screening_num_students - 1:
                        st.session_state.current_student_index += 1
                        st.rerun()
                    else:
                        st.session_state.screening_complete = True
                        st.rerun()
    else:
        results = calculate_screening_results()
        if results:
            st.success("🎉 Screening Complete!")
            st.markdown("---")
            st.subheader("📊 Class Overview")
            total = results["total_students"]
            on_track_pct = (len(results["risk_levels"]["on_track"]) / total) * 100
            monitor_pct = (len(results["risk_levels"]["monitor"]) / total) * 100
            priority_pct = (len(results["risk_levels"]["priority"]) / total) * 100
            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("🟢 On Track", f"{len(results['risk_levels']['on_track'])} ({on_track_pct:.0f}%)")
            with col2:
                st.metric("🟡 Monitor", f"{len(results['risk_levels']['monitor'])} ({monitor_pct:.0f}%)")
            with col3:
                st.metric("🔴 Priority", f"{len(results['risk_levels']['priority'])} ({priority_pct:.0f}%)")

            st.markdown("---")
            st.subheader("📈 Class Competency Breakdown")
            for comp, avg in results["class_averages"].items():
                pct = (avg / 4.0) * 100
                color = "🟢" if avg >= 3.0 else ("🟡" if avg >= 2.5 else "🔴")
                st.markdown(f"**{color} {comp}**: {avg:.1f}/4.0")
                st.progress(pct / 100)
                st.markdown("")

            lowest_comp = min(results["class_averages"], key=results["class_averages"].get)
            st.info(f"**Class Focus Area:** {lowest_comp} - Consider whole-class intervention")
            if st.button("💡 Get Whole-Class Strategies"):
                with st.spinner("Generating personalized class strategies..."):
                    prompt = get_class_strategies_prompt(results, st.session_state.screening_grade)
                    response = call_claude(prompt, max_tokens=3000, stream=False)
                    if response:
                        st.session_state.screening_interventions["class"] = response
                        st.rerun()

            if "class" in st.session_state.screening_interventions:
                st.markdown("---")
                st.markdown(st.session_state.screening_interventions["class"])

            st.markdown("---")
            st.subheader("👥 Students Needing Support")
            if results["risk_levels"]["priority"]:
                st.markdown("### 🔴 Priority Support (Multiple Concerns)")
                for student_id in results["risk_levels"]["priority"]:
                    with st.expander(f"**{student_id}** - Average: {results['students'][student_id]['average']:.1f}/4.0"):
                        student_results = results["students"][student_id]
                        competencies = ["Self-Awareness", "Self-Management", "Social Awareness", "Relationship Skills", "Decision-Making"]
                        for i, comp in enumerate(competencies):
                            score = student_results["scores"][i]
                            color = "🟢" if score >= 3.0 else ("🟡" if score >= 2.5 else "🔴")
                            st.markdown(f"{color} **{comp}**: {score}/4")
                        if st.button(f"🎯 Generate Intervention Plan", key=f"intervention_{student_id}"):
                            with st.spinner("Creating personalized intervention plan..."):
                                prompt = get_intervention_prompt(student_id, student_results, st.session_state.screening_grade)
                                response = call_claude(prompt, max_tokens=2500, stream=False)
                                if response:
                                    st.session_state.screening_interventions[student_id] = response
                                    st.rerun()
                        if student_id in st.session_state.screening_interventions:
                            st.markdown("---")
                            st.markdown(st.session_state.screening_interventions[student_id])

            if results["risk_levels"]["monitor"]:
                st.markdown("### 🟡 Monitor (1-2 Concerns)")
                for student_id in results["risk_levels"]["monitor"]:
                    with st.expander(f"**{student_id}** - Average: {results['students'][student_id]['average']:.1f}/4.0"):
                        student_results = results["students"][student_id]
                        competencies = ["Self-Awareness", "Self-Management", "Social Awareness", "Relationship Skills", "Decision-Making"]
                        for i, comp in enumerate(competencies):
                            score = student_results["scores"][i]
                            color = "🟢" if score >= 3.0 else ("🟡" if score >= 2.5 else "🔴")
                            st.markdown(f"{color} **{comp}**: {score}/4")
                        if st.button(f"🎯 Generate Intervention Plan", key=f"intervention_{student_id}"):
                            with st.spinner("Creating personalized intervention plan..."):
                                prompt = get_intervention_prompt(student_id, student_results, st.session_state.screening_grade)
                                response = call_claude(prompt, max_tokens=2500, stream=False)
                                if response:
                                    st.session_state.screening_interventions[student_id] = response
                                    st.rerun()
                        if student_id in st.session_state.screening_interventions:
                            st.markdown("---")
                            st.markdown(st.session_state.screening_interventions[student_id])

            if results["risk_levels"]["on_track"]:
                with st.expander(f"🟢 Students On Track ({len(results['risk_levels']['on_track'])} students)"):
                    for student_id in results["risk_levels"]["on_track"]:
                        st.markdown(f"✓ **{student_id}** - Average: {results['students'][student_id]['average']:.1f}/4.0")

            st.markdown("---")
            col1, col2 = st.columns(2)
            with col1:
                if st.button("🔄 Start New Screening"):
                    st.session_state.screening_data = {}
                    st.session_state.screening_interventions = {}
                    st.session_state.current_student_index = 0
                    st.session_state.screening_complete = False
                    st.rerun()
            with col2:
                pass

            st.markdown("---")
            st.subheader("📥 Download Assessment Reports")
            col_dl1, col_dl2, col_dl3 = st.columns(3)
            with col_dl1:
                screening_json = save_screening_data()
                if screening_json:
                    st.download_button(
                        label="💾 Save Data (Reload Later)",
                        data=screening_json,
                        file_name=f"sel_screening_data_{datetime.now().strftime('%Y%m%d_%H%M')}.json",
                        mime="application/json",
                        help="Save this file to reload your screening data later"
                    )
            with col_dl2:
                full_report = create_comprehensive_report()
                if full_report:
                    st.download_button(
                        label="📄 Full Report (Text)",
                        data=full_report,
                        file_name=f"sel_screening_report_{datetime.now().strftime('%Y%m%d')}.txt",
                        mime="text/plain",
                        help="Complete report with all intervention plans"
                    )
            with col_dl3:
                full_report = full_report or create_comprehensive_report()
                if full_report:
                    docx_report = create_docx(full_report)
                    st.download_button(
                        label="📝 Full Report (Word)",
                        data=docx_report,
                        file_name=f"sel_screening_report_{datetime.now().strftime('%Y%m%d')}.docx",
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        help="Professional Word document with all assessments"
                    )
            with st.expander("ℹ️ What's included in downloads"):
                st.markdown("""
                **Data File (.json):** 
                - All screening scores
                - Can be reloaded for future editing
                
                **Full Report (.txt/.docx):**
                - Class overview and statistics
                - Competency breakdowns
                - Whole-class strategies (if generated)
                - Individual intervention plans (if generated)
                - Student groupings by support level
                """)

# ---- COMMON OUTPUT AREA (Tabs 1 & 2) ----
if st.session_state.ai_response:
    st.markdown("---")
    st.header(st.session_state.response_title)
    st.markdown(st.session_state.ai_response)

    st.markdown("---")
    st.subheader("📧 Parent Communication")
    if st.button("Generate Parent Email"):
        with st.spinner("Drafting a parent email..."):
            email_prompt = get_parent_email_prompt(st.session_state.ai_response)
            response = call_claude(email_prompt, max_tokens=2048, stream=False)
            if response:
                st.session_state.parent_email = response
    if st.session_state.parent_email:
        st.text_area("Parent Email Draft", value=st.session_state.parent_email, height=300)

    st.markdown("---")
    st.subheader("👩‍🏫 Generate Student-Facing Materials")
    if st.button("Generate Materials"):
        with st.spinner("✍️ Creating student materials..."):
            materials_prompt = get_student_materials_prompt(st.session_state.ai_response)
            response = call_claude(materials_prompt, stream=False)
            if response:
                st.session_state.student_materials = response
    if st.session_state.student_materials:
        st.markdown(st.session_state.student_materials)

    st.markdown("---")
    st.subheader("🧠 Differentiate This Lesson")
    if st.button("Generate Differentiation Strategies"):
        with st.spinner("💡 Coming up with strategies for diverse learners..."):
            diff_prompt = get_differentiation_prompt(st.session_state.ai_response)
            response = call_claude(diff_prompt, stream=False)
            if response:
                st.session_state.differentiation_response = response
    if st.session_state.differentiation_response:
        st.markdown(st.session_state.differentiation_response)

    st.markdown("---")
    st.subheader("📥 Download Your Plan")
    full_download_text = st.session_state.ai_response
    if st.session_state.parent_email:
        full_download_text += "\n\n---\n\n# Parent Communication Draft\n\n" + st.session_state.parent_email
    if st.session_state.student_materials:
        full_download_text += "\n\n---\n\n# Student-Facing Materials\n\n" + st.session_state.student_materials
    if st.session_state.differentiation_response:
        full_download_text += "\n\n---\n\n# Differentiation Strategies\n\n" + st.session_state.differentiation_response
    if full_download_text.strip():
        docx_file = create_docx(full_download_text)
        dl_col1, dl_col2 = st.columns(2)
        with dl_col1:
            st.download_button(label="Download as Text File (.txt)", data=full_download_text.encode('utf-8-sig'), file_name="sel_plan.txt", mime="text/plain")
        with dl_col2:
            if docx_file:
                st.download_button(label="Download as Word Doc (.docx)", data=docx_file, file_name="sel_plan.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

st.markdown("---")
st.markdown("*💡 Powered by Claude Sonnet 4.5 from Anthropic*")
st.caption(f"Session started: {st.session_state.session_start_time.strftime('%Y-%m-%d %H:%M:%S')}")

